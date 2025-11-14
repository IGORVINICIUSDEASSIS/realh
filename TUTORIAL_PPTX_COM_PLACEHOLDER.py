"""
TUTORIAL COMPLETO: Como criar um PPTX com Placeholders
Documentação em Português - 100% prático
"""

# ============================================================================
# PARTE 1: O QUE É UM PLACEHOLDER?
# ============================================================================
"""
Um placeholder é um "espaço reservado" que você coloca no PowerPoint.
Quando o Python executa, ele SUBSTITUI esse espaço pelo valor real.

EXEMPLO:
--------
Você digita no PowerPoint:  "Relatório {{TITULO}} - {{PERIODO}}"
Python substitui:          "Relatório Vendas - Nov 2024"

É tipo um "template" de carta, sabe? Você cria a carta UMA VEZ,
depois só muda o nome e a data! 📝
"""

# ============================================================================
# PARTE 2: COMO CRIAR UM PPTX COM PLACEHOLDER NO POWERPOINT
# ============================================================================
"""
PASSO 1: Abrir PowerPoint
─────────────────────────
1. Abra PowerPoint (ou use online.office.com)
2. Crie uma apresentação nova ou abra a existente

PASSO 2: Adicionar Placeholder de Texto
───────────────────────────────────────
1. Clique em "Inserir" → "Caixa de Texto"
2. Desenhe uma caixa no slide
3. Digite: {{TITULO}}
4. Formate como quiser (fonte, tamanho, cor)

PASSO 3: Salvar
───────────────
1. Arquivo → Salvar como
2. Nome: "template_relatorio.pptx"
3. Formato: "Apresentação do PowerPoint (*.pptx)"
4. Local: Na mesma pasta do projeto Python

RESULTADO:
──────────
Seu PowerPoint está pronto com placeholders!
O Python vai ler e substituir os {{VALORES}}
"""

# ============================================================================
# PARTE 3: COMO O PYTHON LÊ E PREENCHE
# ============================================================================
"""
Quando você executa Python, ele faz isto:

1. ABRE o arquivo PPTX
   └─ prs = Presentation('template_relatorio.pptx')

2. VAI SLIDE POR SLIDE
   └─ for slide in prs.slides:

3. VAI CAIXA DE TEXTO POR CAIXA DE TEXTO
   └─ for shape in slide.shapes:
      if hasattr(shape, 'text_frame'):

4. PROCURA POR {{PLACEHOLDER}}
   └─ if '{{TITULO}}' in run.text:

5. SUBSTITUI PELO VALOR REAL
   └─ run.text = run.text.replace('{{TITULO}}', 'Vendas')

6. SALVA NOVO ARQUIVO
   └─ prs.save('relatorio_final.pptx')

BINGO! 🎉
"""

# ============================================================================
# PARTE 4: EXEMPLOS PRÁTICOS DE PLACEHOLDERS
# ============================================================================

# Exemplo 1: TEXTO SIMPLES
# ──────────────────────────
# No PowerPoint escreva:
#   "Relatório de {{TIPO}} - {{PERIODO}}"
#
# No Python:
#   substituicoes = {
#       '{{TIPO}}': 'Vendas',
#       '{{PERIODO}}': 'Novembro 2024'
#   }
#   # Resultado: "Relatório de Vendas - Novembro 2024"


# Exemplo 2: VALORES COM FORMATAÇÃO
# ──────────────────────────────────
# No PowerPoint escreva:
#   "Total: R$ {{VENDAS_TOTAIS}}"
#
# No Python:
#   substituicoes = {
#       '{{VENDAS_TOTAIS}}': '1.500.000,00'
#   }
#   # Resultado: "Total: R$ 1.500.000,00"


# Exemplo 3: MÚLTIPLAS LINHAS (Métricas)
# ───────────────────────────────────────
# No PowerPoint escreva:
#   "{{METRICAS}}"
#
# No Python:
#   metricas_texto = """• Vendas: R$ 1.5M
# • Clientes: 500
# • Produto Top: XYZ"""
#   substituicoes = {
#       '{{METRICAS}}': metricas_texto
#   }


# Exemplo 4: IMAGENS (Gráficos)
# ──────────────────────────────
# No PowerPoint escreva:
#   "[{{GRAFICO}}]"
#
# No Python:
#   # REMOVE o texto placeholder
#   # INSERE a imagem do gráfico no lugar
#   # (Isso é mais complexo, veja próxima seção)


# ============================================================================
# PARTE 5: CÓDIGO PYTHON QUE FAZ SUBSTITUIÇÃO
# ============================================================================

def exemplo_substituicao_simples():
    """Exemplo básico de como o Python substitui placeholders"""
    
    from pptx import Presentation
    
    # PASSO 1: Abrir o template
    prs = Presentation('template_relatorio.pptx')
    
    # PASSO 2: Definir o que substituir
    substituicoes = {
        '{{TITULO}}': 'Relatório de Vendas',
        '{{PERIODO}}': 'Novembro 2024',
        '{{VENDAS}}': 'R$ 1.500.000,00',
        '{{CLIENTES}}': '500 clientes',
    }
    
    # PASSO 3: Percorrer TODOS os slides
    for slide in prs.slides:
        # PASSO 4: Percorrer TODAS as formas (caixas de texto)
        for shape in slide.shapes:
            # PASSO 5: Verificar se é uma caixa de texto
            if hasattr(shape, 'text_frame'):
                # PASSO 6: Percorrer cada parágrafo
                for paragraph in shape.text_frame.paragraphs:
                    # PASSO 7: Percorrer cada "run" (pedaço de texto)
                    for run in paragraph.runs:
                        # PASSO 8: Fazer substituição
                        for placeholder, valor in substituicoes.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(valor))
    
    # PASSO 9: Salvar arquivo preenchido
    prs.save('relatorio_preenchido.pptx')
    print("✅ Relatório gerado: relatorio_preenchido.pptx")


# ============================================================================
# PARTE 6: PLACEHOLDERS QUE JÁ TEMOS NO SEU SISTEMA
# ============================================================================

"""
No seu projeto, já implementamos estes placeholders:

✅ {{TITULO}}          → Título do relatório
✅ {{PERIODO}}         → Período (ex: Nov/2024)
✅ {{METRICAS}}        → Lista de métricas principais
✅ {{GRAFICO}}         → Imagem dos gráficos (como PNG)
✅ {{NOME_GRAFICO}}    → Nome do tipo de gráfico

COMO USAR:
──────────
1. Crie um PPTX normal no PowerPoint
2. Adicione caixas de texto com esses placeholders
3. Formate como quiser (cores, fonts, tamanho)
4. Salve como: template_relatorio.pptx
5. O Python automaticamente:
   - Lê o arquivo
   - Procura pelos placeholders
   - Substitui pelos dados reais
   - Salva um novo arquivo com tudo preenchido
   - Você baixa o novo arquivo! 📥
"""

# ============================================================================
# PARTE 7: PASSO-A-PASSO PRÁTICO (COM IMAGENS)
# ============================================================================

"""
TUTORIAL VISUAL - COMO FAZER NO POWERPOINT:
==============================================

┌─────────────────────────────────────────────┐
│ PASSO 1: Abra PowerPoint                    │
├─────────────────────────────────────────────┤
│ 1. Abra PowerPoint ou vai pra office.com   │
│ 2. Clique em "Apresentação em Branco"      │
└─────────────────────────────────────────────┘

┌─────────────────────────────────────────────┐
│ PASSO 2: Insira uma Caixa de Texto          │
├─────────────────────────────────────────────┤
│ 1. Clique em "Inserir" no menu superior    │
│ 2. Procure por "Caixa de Texto"            │
│ 3. Desenhe uma caixa no slide              │
│    (clique e arraste)                      │
└─────────────────────────────────────────────┘

┌─────────────────────────────────────────────┐
│ PASSO 3: Digite o Placeholder               │
├─────────────────────────────────────────────┤
│ 1. Dentro da caixa, digite:                │
│    {{TITULO}}                              │
│    {{PERIODO}}                             │
│    {{METRICAS}}                            │
│                                             │
│ ⚠️ Importante: Escreva EXATAMENTE assim!    │
│    Com as chaves duplas {{ }}               │
└─────────────────────────────────────────────┘

┌─────────────────────────────────────────────┐
│ PASSO 4: Formate (Opcional)                │
├─────────────────────────────────────────────┤
│ 1. Selecione o texto                       │
│ 2. Altere:                                 │
│    • Fonte (Arial, Calibri, etc)           │
│    • Tamanho (18pt, 24pt, etc)            │
│    • Cor (verde, azul, etc)                │
│    • Alinhamento (centro, esquerda, etc)   │
└─────────────────────────────────────────────┘

┌─────────────────────────────────────────────┐
│ PASSO 5: Salve o Arquivo                   │
├─────────────────────────────────────────────┤
│ 1. Clique em "Arquivo"                     │
│ 2. Clique em "Salvar como"                 │
│ 3. Nome: template_relatorio.pptx          │
│ 4. Formato: PowerPoint (*.pptx)            │
│ 5. Local: /workspaces/realh/              │
│ 6. Clique "Salvar"                         │
└─────────────────────────────────────────────┘

PRONTO! 🎉
Seu template está criado com placeholders!
"""

# ============================================================================
# PARTE 8: DICAS E TRUQUES
# ============================================================================

"""
DICA 1: Múltiplos Placeholders no Mesmo Slide
──────────────────────────────────────────────
Você pode colocar vários {{PLACEHOLDER}} em um mesmo slide.
Python vai substituir TODOS automaticamente.

Exemplo no PowerPoint:
  ┌──────────────────────────┐
  │ {{TITULO}}              │
  │ Período: {{PERIODO}}    │
  │ Vendas: {{VENDAS}}      │
  └──────────────────────────┘

Python vai substituir os 3!


DICA 2: Preservar Formatação
─────────────────────────────
Se você quer que o texto fique VERMELHO mesmo após substituição:

No PowerPoint:
  1. Escreva: {{VENDAS}}
  2. Selecione o texto
  3. Mude para VERMELHO
  4. Python mantém a cor vermelha! ✓


DICA 3: Placeholders em Tabelas
────────────────────────────────
Você pode colocar placeholders DENTRO de tabelas!

No PowerPoint:
  1. Insira uma tabela (Inserir → Tabela)
  2. Em uma célula, escreva: {{CLIENTE}}
  3. Python substitui! ✓


DICA 4: Quebras de Linha em Métricas
──────────────────────────────────────
Para múltiplas métricas ficarem em linhas diferentes:

No Python:
  metricas = '''• Vendas: R$ 1M
• Clientes: 500
• Produtos: 50'''
  
No PowerPoint (será exibido assim):
  • Vendas: R$ 1M
  • Clientes: 500
  • Produtos: 50


DICA 5: Testar o Template
──────────────────────────
1. Crie um template simples primeiro
2. Teste com dados de exemplo
3. Se funcionar, customize o design
4. Deixe os placeholders {{INTACTOS}}
"""

# ============================================================================
# PARTE 9: ERROS COMUNS (E COMO EVITAR)
# ============================================================================

"""
❌ ERRO 1: Escrever diferente
──────────────────────────────
Errado: {{titulo}}  ou  {{TITULO }}  ou  {TITULO}
Certo:  {{TITULO}}

Python procura EXATAMENTE por "{{TITULO}}"
Se escrever diferente, não encontra!


❌ ERRO 2: Confundir maiúsculas e minúsculas
──────────────────────────────────────────────
Errado:
  No PowerPoint: {{Titulo}}
  No Python: '{{TITULO}}'
  → NÃO vai funcionar!

Certo: Sempre use maiúsculas em ambos


❌ ERRO 3: Deletar o placeholder por acaso
────────────────────────────────────────────
Cuidado: Quando você edita o PowerPoint depois,
não delete os {{PLACEHOLDERS}} acidentalmente!


❌ ERRO 4: Salvar como .pptm ou .odp
──────────────────────────────────────
Errado: .pptm (macro) ou .odp (OpenOffice)
Certo:  .pptx (PowerPoint moderno)


❌ ERRO 5: Colocar placeholder no slide master
────────────────────────────────────────────────
Coloque em slides NORMAIS, não no master!
"""

# ============================================================================
# PARTE 10: RESUMO EM UMA PÁGINA
# ============================================================================

"""
╔═══════════════════════════════════════════════════════════════════╗
║              RESUMO RÁPIDO - CRIAR PPTX COM PLACEHOLDER           ║
╚═══════════════════════════════════════════════════════════════════╝

1. NO POWERPOINT:
   ✓ Abra PowerPoint
   ✓ Insira Caixa de Texto (Inserir → Caixa de Texto)
   ✓ Escreva: {{TITULO}}, {{PERIODO}}, {{METRICAS}}, etc
   ✓ Formate como quiser (cores, fonts, tamanho)
   ✓ Salve como: template_relatorio.pptx

2. NO PYTHON:
   ✓ Python lê o arquivo PPTX
   ✓ Procura por {{PLACEHOLDER}}
   ✓ Substitui pelos valores reais
   ✓ Salva novo PPTX preenchido
   ✓ Você baixa! ✅

3. PLACEHOLDERS DISPONÍVEIS:
   {{TITULO}}         → Título do relatório
   {{PERIODO}}        → Período
   {{METRICAS}}       → Métricas em lista
   {{GRAFICO}}        → Imagem do gráfico
   {{NOME_GRAFICO}}   → Nome do gráfico

4. ABRA O ARQUIVO:
   /workspaces/realh/template_relatorio.pptx

5. CUSTOMIZE:
   Desenhe, escreva, decore como quiser!
   Só não apague os {{PLACEHOLDERS}}

6. SALVE COMO:
   template_relatorio_customizado.pptx

7. VOLTE AO APP:
   📄 Relatório → 📋 Usar Template
   Selecione seu arquivo
   Clique "Gerar"
   Baixe o relatório! 🎉
"""

# ============================================================================
# LINKS E REFERÊNCIAS ÚTEIS
# ============================================================================

"""
LINKS E REFERÊNCIAS:
═══════════════════

Se quiser aprender mais:

1. DOCUMENTAÇÃO PYTHON-PPTX:
   https://python-pptx.readthedocs.io/
   (Em inglês, mas é a referência oficial)

2. TUTORIAL VISUAL (YouTube):
   Procure por "python-pptx tutorial"
   Tem vários em inglês que ajudam

3. OFICINA MICROSOFT OFFICE:
   https://www.microsoft.com/pt-br/office
   (Para aprender PowerPoint mesmo)

4. STACK OVERFLOW:
   https://stackoverflow.com/questions/tagged/python-pptx
   (Comunidade que responde dúvidas)

MAS AQUI NO SEU CÓDIGO:
─────────────────────
Você já tem tudo pronto!
Só precisa:
  1. Abrir /workspaces/realh/template_relatorio.pptx
  2. Editar no PowerPoint
  3. Deixar os {{PLACEHOLDERS}}
  4. Salvar
  5. Python faz o resto! 🤖
"""

print(__doc__)
