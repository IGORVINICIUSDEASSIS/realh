"""
Funções para trabalhar com templates de PPTX
Permite usar um arquivo PowerPoint como base e preencher com dados
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import re
import os
import tempfile
import plotly.io as pio


def gerar_template_padrao(caminho_template='template_relatorio.pptx'):
    """
    Gera um arquivo template PPTX padrão que você pode customizar no PowerPoint.
    
    O template contém placeholders que o Python vai preencher:
    {{TITULO}} - Título do relatório
    {{PERIODO}} - Período
    {{VENDAS_TOTAIS}} - Métrica de vendas
    {{GRAFICO_CLIENTES}} - Espaço para gráfico
    
    Args:
        caminho_template: Caminho onde salvar o template
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    cor_primaria = RGBColor(0, 204, 150)
    cor_titulo = RGBColor(33, 37, 41)
    cor_texto = RGBColor(50, 50, 50)
    
    # ===== SLIDE 1: CAPA =====
    slide_capa = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide_capa.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = cor_titulo
    
    # Logo
    logo_path = 'assets/logotipo.png'
    if os.path.exists(logo_path):
        try:
            slide_capa.shapes.add_picture(logo_path, Inches(3.5), Inches(0.5), height=Inches(1.2))
        except:
            pass
    
    # Placeholder: Título
    txBox_titulo = slide_capa.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf_titulo = txBox_titulo.text_frame
    tf_titulo.word_wrap = True
    p_titulo = tf_titulo.paragraphs[0]
    p_titulo.text = "{{TITULO}}"
    p_titulo.font.size = Pt(54)
    p_titulo.font.bold = True
    p_titulo.font.color.rgb = cor_primaria
    p_titulo.alignment = PP_ALIGN.CENTER
    
    # Placeholder: Período
    txBox_periodo = slide_capa.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf_periodo = txBox_periodo.text_frame
    p_periodo = tf_periodo.paragraphs[0]
    p_periodo.text = "Período: {{PERIODO}}"
    p_periodo.font.size = Pt(28)
    p_periodo.font.color.rgb = RGBColor(200, 200, 200)
    p_periodo.alignment = PP_ALIGN.CENTER
    
    # ===== SLIDE 2: MÉTRICAS =====
    slide_metricas = prs.slides.add_slide(prs.slide_layouts[6])
    background2 = slide_metricas.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Logo
    if os.path.exists(logo_path):
        try:
            slide_metricas.shapes.add_picture(logo_path, Inches(0.3), Inches(0.1), height=Inches(0.6))
        except:
            pass
    
    # Título
    txBox_met_titulo = slide_metricas.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf_met_titulo = txBox_met_titulo.text_frame
    p_met_titulo = tf_met_titulo.paragraphs[0]
    p_met_titulo.text = "📊 Métricas Principais"
    p_met_titulo.font.size = Pt(40)
    p_met_titulo.font.bold = True
    p_met_titulo.font.color.rgb = cor_titulo
    
    # Placeholder: Métricas
    txBox_metricas = slide_metricas.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    tf_metricas = txBox_metricas.text_frame
    tf_metricas.word_wrap = True
    p_met = tf_metricas.paragraphs[0]
    p_met.text = "{{METRICAS}}"
    p_met.font.size = Pt(18)
    p_met.font.color.rgb = cor_texto
    
    # ===== SLIDE 3: GRÁFICO =====
    slide_grafico = prs.slides.add_slide(prs.slide_layouts[6])
    background3 = slide_grafico.background
    fill3 = background3.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Logo
    if os.path.exists(logo_path):
        try:
            slide_grafico.shapes.add_picture(logo_path, Inches(0.3), Inches(0.1), height=Inches(0.6))
        except:
            pass
    
    # Título
    txBox_graf_titulo = slide_grafico.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf_graf_titulo = txBox_graf_titulo.text_frame
    p_graf_titulo = tf_graf_titulo.paragraphs[0]
    p_graf_titulo.text = "📊 Gráfico - {{NOME_GRAFICO}}"
    p_graf_titulo.font.size = Pt(36)
    p_graf_titulo.font.bold = True
    p_graf_titulo.font.color.rgb = cor_titulo
    
    # Placeholder: Imagem do gráfico (deixa espaço reservado)
    # Nota: {{GRAFICO_CLIENTES}} será substituído por uma imagem
    txBox_placeholder_grafico = slide_grafico.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(5.5))
    tf_placeholder_grafico = txBox_placeholder_grafico.text_frame
    p_placeholder_grafico = tf_placeholder_grafico.paragraphs[0]
    p_placeholder_grafico.text = "[{{GRAFICO}}]"
    p_placeholder_grafico.font.size = Pt(16)
    p_placeholder_grafico.font.color.rgb = RGBColor(200, 200, 200)
    p_placeholder_grafico.alignment = PP_ALIGN.CENTER
    
    # Salvar template
    prs.save(caminho_template)
    print(f"✅ Template criado: {caminho_template}")
    print(f"   Você pode agora abrir no PowerPoint e customizar o design!")
    print(f"\n   Placeholders disponíveis:")
    print(f"   - {{{{TITULO}}}}")
    print(f"   - {{{{PERIODO}}}}")
    print(f"   - {{{{METRICAS}}}}")
    print(f"   - {{{{GRAFICO}}}}")
    print(f"   - {{{{NOME_GRAFICO}}}}")


def preencher_template_pptx(caminho_template, titulo, periodo, metricas_dict, graficos_dict=None):
    """
    Lê um template PPTX e preenche os placeholders com dados reais.
    
    Args:
        caminho_template: Caminho do arquivo template.pptx
        titulo: Título do relatório
        periodo: Período (ex: "Set/2024")
        metricas_dict: Dict com métricas {'Vendas Totais': 'R$ 10.000'}
        graficos_dict: Dict com figuras Plotly {'Nome': fig, ...}
    
    Returns:
        bytes: Arquivo PPTX preenchido em bytes
    """
    
    prs = Presentation(caminho_template)
    logo_path = 'assets/logotipo.png'
    
    # Dicionário de substituições
    substituicoes = {
        '{{TITULO}}': titulo,
        '{{PERIODO}}': periodo,
    }
    
    # Montar string de métricas
    metricas_texto = ""
    for nome, valor in metricas_dict.items():
        metricas_texto += f"• {nome}: {valor}\n"
    substituicoes['{{METRICAS}}'] = metricas_texto.strip()
    
    # Varrer todos os slides
    for slide_idx, slide in enumerate(prs.slides):
        # Procurar texto nos shapes
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        texto_original = run.text
                        
                        # Substituir placeholders
                        for placeholder, valor in substituicoes.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(valor))
    
    # ===== INSERIR GRÁFICOS =====
    if graficos_dict:
        for titulo_grafico, fig in graficos_dict.items():
            try:
                # Converter gráfico para imagem
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    tmp_path = tmp.name
                
                # Tentar usar kaleido primeiro
                try:
                    pio.write_image(fig, tmp_path, width=1200, height=700, format='png', engine='kaleido')
                except Exception:
                    # Se kaleido falhar, tentar com outros engines
                    try:
                        pio.write_image(fig, tmp_path, width=1200, height=700, format='png', engine='orca')
                    except Exception:
                        # Se tudo falhar, pular este gráfico
                        print(f"ℹ️ Gráfico '{titulo_grafico}' não pôde ser convertido - requer Chrome/Chromium")
                        continue
                
                if os.path.exists(tmp_path) and os.path.getsize(tmp_path) > 0:
                    # Procurar slide com placeholder de gráfico e substituir
                    for slide_idx, slide in enumerate(prs.slides):
                        for shape_idx, shape in enumerate(slide.shapes):
                            if hasattr(shape, 'text_frame') and '{{GRAFICO}}' in shape.text:
                                # Remover o placeholder de texto
                                p = shape.text_frame.paragraphs[0]
                                p.text = ""
                                
                                # Adicionar imagem do gráfico
                                slide.shapes.add_picture(
                                    tmp_path, 
                                    Inches(0.5), 
                                    Inches(1.1), 
                                    width=Inches(9)
                                )
                                break
                    
                    os.unlink(tmp_path)
            except Exception as e:
                print(f"ℹ️ Gráfico '{titulo_grafico}' não incluído - requer Chrome/Chromium instalado")
    
    # Salvar em bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()
