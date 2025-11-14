"""
Funções utilitárias compartilhadas entre as páginas do dashboard
"""
import pandas as pd
import streamlit as st
from dateutil.relativedelta import relativedelta
import os

# ==============================
# FUNÇÕES DE FORMATAÇÃO SEGURA
# ==============================
def safe_strftime(date_value, format_str='%d/%m/%Y'):
    """
    Formata uma data de forma segura, tratando valores NaT/None
    
    Args:
        date_value: datetime object ou NaT
        format_str: formato de saída (padrão: '%d/%m/%Y')
    
    Returns:
        str: data formatada ou mensagem de erro
    """
    if pd.isna(date_value):
        return "Data inválida"
    try:
        return date_value.strftime(format_str)
    except:
        return "Data inválida"

# ==============================
# FUNÇÕES DE MÊS COMERCIAL
# ==============================
def calcular_mes_comercial(data):
    """
    Calcula o mês comercial baseado na regra: 16/MM ao 15/MM+1
    Exemplo: 16/09/2024 até 15/10/2024 = "Set/2024"
    
    Args:
        data: datetime object
    
    Returns:
        str: Nome do mês comercial no formato "MMM/YYYY"
    """
    if pd.isna(data):
        return None
    
    # Se o dia é >= 16, pertence ao mês comercial atual
    if data.day >= 16:
        mes_comercial = data
    else:
        # Se o dia é < 16, pertence ao mês comercial anterior
        mes_comercial = data - relativedelta(months=1)
    
    # Retornar no formato "Set/2024"
    meses_pt = {
        1: 'Jan', 2: 'Fev', 3: 'Mar', 4: 'Abr', 5: 'Mai', 6: 'Jun',
        7: 'Jul', 8: 'Ago', 9: 'Set', 10: 'Out', 11: 'Nov', 12: 'Dez'
    }
    
    return f"{meses_pt[mes_comercial.month]}/{mes_comercial.year}"

def obter_periodo_mes_comercial(mes_comercial_str):
    """
    Retorna as datas de início e fim de um mês comercial
    
    Args:
        mes_comercial_str: str no formato "MMM/YYYY" (ex: "Set/2024")
    
    Returns:
        tuple: (data_inicio, data_fim)
    """
    meses_pt_inv = {
        'Jan': 1, 'Fev': 2, 'Mar': 3, 'Abr': 4, 'Mai': 5, 'Jun': 6,
        'Jul': 7, 'Ago': 8, 'Set': 9, 'Out': 10, 'Nov': 11, 'Dez': 12
    }
    
    mes_str, ano_str = mes_comercial_str.split('/')
    mes = meses_pt_inv[mes_str]
    ano = int(ano_str)
    
    # Início: dia 16 do mês
    data_inicio = pd.Timestamp(year=ano, month=mes, day=16)
    
    # Fim: dia 15 do mês seguinte
    if mes == 12:
        data_fim = pd.Timestamp(year=ano + 1, month=1, day=15, hour=23, minute=59, second=59)
    else:
        data_fim = pd.Timestamp(year=ano, month=mes + 1, day=15, hour=23, minute=59, second=59)
    
    return data_inicio, data_fim

def ordenar_mes_comercial(mes_str):
    """Converte mês comercial em timestamp para ordenação"""
    meses_pt_inv = {
        'Jan': 1, 'Fev': 2, 'Mar': 3, 'Abr': 4, 'Mai': 5, 'Jun': 6,
        'Jul': 7, 'Ago': 8, 'Set': 9, 'Out': 10, 'Nov': 11, 'Dez': 12
    }
    mes, ano = mes_str.split('/')
    return pd.Timestamp(year=int(ano), month=meses_pt_inv[mes], day=1)

def obter_mes_comercial_atual():
    """Retorna o mês comercial atual (mesmo que incompleto)"""
    return calcular_mes_comercial(pd.Timestamp.now())

def mes_comercial_esta_completo(mes_comercial_str):
    """
    Verifica se um mês comercial já está completo (terminou no dia 15).
    Retorna True se já passou do dia 15 do mês seguinte.
    """
    data_inicio, data_fim = obter_periodo_mes_comercial(mes_comercial_str)
    hoje = pd.Timestamp.now()
    return hoje > data_fim

# ==============================
# FUNÇÕES DE FORMATAÇÃO
# ==============================
def formatar_moeda(valor):
    """Formata valor em moeda brasileira"""
    try:
        return f"R$ {valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    except:
        return "R$ 0,00"

def detectar_coluna_hierarquica(df_columns, nomes_possiveis):
    """
    Detecta automaticamente colunas com base em nomes similares.
    Retorna o índice da coluna encontrada + 1 (para compensar "Nenhuma") ou 0 se não encontrar.
    """
    for nome in nomes_possiveis:
        for idx, col in enumerate(df_columns):
            if nome.lower() in col.lower():
                return idx + 1  # +1 porque "Nenhuma" está na posição 0
    return 0

# ==============================
# FUNÇÕES DE UI
# ==============================
def exibir_top_com_alternancia(df, titulo, chave_session, tipo_grafico='bar'):
    """
    Exibe um gráfico por padrão e oferece opção de alternar para tabela com TODOS os dados.
    
    Args:
        df: DataFrame com os dados já ordenados
        titulo: Título da seção
        chave_session: Chave única para armazenar estado no session_state
        tipo_grafico: 'bar' (padrão), 'horizontal_bar', 'pie', etc
    
    Returns:
        None (exibe na tela)
    """
    import plotly.graph_objects as go
    
    if f"{chave_session}_modo_tabela" not in st.session_state:
        st.session_state[f"{chave_session}_modo_tabela"] = False
    
    st.markdown(f"#### {titulo}")
    
    # Botão para alternar entre gráfico e tabela
    col_btn1, col_btn2 = st.columns([0.1, 0.9])
    with col_btn1:
        if st.button(
            "📊" if st.session_state[f"{chave_session}_modo_tabela"] else "📋",
            key=f"{chave_session}_btn",
            help="Alternar entre gráfico e tabela"
        ):
            st.session_state[f"{chave_session}_modo_tabela"] = not st.session_state[f"{chave_session}_modo_tabela"]
            st.rerun()
    
    with col_btn2:
        if st.session_state[f"{chave_session}_modo_tabela"]:
            st.markdown("**Exibindo: Tabela Completa**")
        else:
            st.markdown("**Exibindo: Gráfico**")
    
    st.markdown("")
    
    if st.session_state[f"{chave_session}_modo_tabela"]:
        # Modo tabela - exibe TODOS os dados
        st.dataframe(df, use_container_width=True, hide_index=True)
    else:
        # Modo gráfico - exibe top 10
        df_top = df.head(10)
        
        if len(df_top) > 0:
            if tipo_grafico == 'bar':
                # Gráfico de barras horizontal com valores visíveis nas barras
                fig = go.Figure()
                fig.add_trace(go.Bar(
                    y=df_top.iloc[:, 0],
                    x=df_top.iloc[:, 1],
                    orientation='h',
                    marker_color='#EF553B',
                    text=df_top.iloc[:, 1],
                    textposition='outside',
                    hovertemplate='<b>%{y}</b><extra></extra>'
                ))
                fig.update_layout(
                    height=400,
                    margin=dict(l=200, r=100, t=0, b=50),
                    yaxis={'categoryorder': 'total ascending'},
                    showlegend=False,
                    xaxis=dict(showticklabels=False, showline=False, zeroline=False),
                    yaxis_title=""
                )
                st.plotly_chart(fig, use_container_width=True)
            
            elif tipo_grafico == 'pie':
                # Gráfico de pizza com valores e percentuais
                fig = go.Figure(data=[go.Pie(
                    labels=df_top.iloc[:, 0],
                    values=df_top.iloc[:, 1],
                    textposition='inside',
                    textinfo='label+percent',
                    hovertemplate='<b>%{label}</b><br>%{value}<br>%{percent}<extra></extra>'
                )])
                fig.update_layout(height=400)
                st.plotly_chart(fig, use_container_width=True)
            
            else:
                # Padrão: tabela dos top 10
                st.dataframe(df_top, use_container_width=True, hide_index=True)

def exibir_logo():
    """
    Exibe o logotipo no canto superior direito da página.
    Procura o arquivo logotipo.png na pasta assets.
    """
    # Exibir informações na sidebar
    with st.sidebar:
        # Exibir usuário logado
        if 'authenticated' in st.session_state and st.session_state['authenticated']:
            user_data = st.session_state.get('user_data', {})
            st.markdown(f"**👤 Usuário:** {user_data.get('nome', 'N/A')}")
            st.markdown(f"**🔑 Tipo:** {user_data.get('tipo', 'N/A').title()}")
            
            # Exibir data/hora do último upload se existir
            config = st.session_state.get('config', {})
            if config and 'data_hora_upload' in config:
                st.markdown("---")
                st.markdown(f"**📅 Última Atualização:**")
                st.markdown(f"*{config['data_hora_upload']}*")
            
            st.markdown("---")
    
    logo_path = os.path.join(os.path.dirname(__file__), "assets", "logotipo.png")
    
    if os.path.exists(logo_path):
        # Converte imagem para base64
        with open(logo_path, "rb") as f:
            import base64
            logo_base64 = base64.b64encode(f.read()).decode()
        
        # CSS para posicionar o logo no canto superior direito
        st.markdown(
            f"""
            <style>
            .logo-container {{
                position: fixed;
                top: 3.5rem;
                right: 2rem;
                z-index: 999;
                background: white;
                padding: 0.5rem;
                border-radius: 8px;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }}
            .logo-container img {{
                width: 150px;
                height: auto;
                display: block;
            }}
            @media (max-width: 768px) {{
                .logo-container {{
                    width: 100px;
                }}
                .logo-container img {{
                    width: 100px;
                }}
            }}
            </style>
            <div class="logo-container">
                <img src="data:image/png;base64,{logo_base64}" alt="Logo">
            </div>
            """,
            unsafe_allow_html=True
        )
    else:
        # Se não encontrar o logo, não faz nada (silencioso)
        pass

# ==============================
# FUNÇÕES PARA GERAÇÃO DE RELATÓRIOS (PPTX)
# ==============================
def gerar_relatorio_pptx(titulo, periodo, metricas_dict, tops_dict, graficos_dict=None):
    """
    Gera um arquivo PPTX com relatório de vendas com gráficos.
    
    Args:
        titulo: Título do relatório
        periodo: Período (ex: "Set/2024")
        metricas_dict: Dict com métricas principais {'nome': valor}
        tops_dict: Dict com DataFrames dos tops {'Top Clientes': df, ...}
        graficos_dict: Dict com figuras Plotly {'Nome': fig, ...}
    
    Returns:
        bytes: Arquivo PPTX em bytes para download
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from io import BytesIO
    import plotly.graph_objects as go
    import tempfile
    import os
    
    # Criar apresentação
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Cores
    cor_primaria = RGBColor(0, 204, 150)  # Verde
    cor_texto = RGBColor(50, 50, 50)
    cor_titulo = RGBColor(33, 37, 41)
    cor_fundo_claro = RGBColor(245, 245, 245)
    
    # ===== SLIDE 1: CAPA =====
    slide_capa = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide_capa.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = cor_titulo
    
    # Adicionar logo no topo
    logo_path = 'assets/logotipo.png'
    if os.path.exists(logo_path):
        try:
            slide_capa.shapes.add_picture(logo_path, Inches(3.5), Inches(0.5), height=Inches(1.2))
        except:
            pass  # Se falhar, continua sem logo
    
    # Título
    txBox = slide_capa.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = cor_primaria
    p.alignment = PP_ALIGN.CENTER
    
    # Período
    txBox2 = slide_capa.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Período: {periodo}"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    
    # ===== SLIDE 2: MÉTRICAS =====
    slide_metricas = prs.slides.add_slide(prs.slide_layouts[6])
    background2 = slide_metricas.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Adicionar logo no topo
    if os.path.exists(logo_path):
        try:
            slide_metricas.shapes.add_picture(logo_path, Inches(0.3), Inches(0.1), height=Inches(0.6))
        except:
            pass
    
    # Título da aba
    txBox_titulo = slide_metricas.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf_titulo = txBox_titulo.text_frame
    p_titulo = tf_titulo.paragraphs[0]
    p_titulo.text = "📊 Métricas Principais"
    p_titulo.font.size = Pt(40)
    p_titulo.font.bold = True
    p_titulo.font.color.rgb = cor_titulo
    
    # Adicionar métricas em grade
    for idx, (nome, valor) in enumerate(metricas_dict.items()):
        linha = idx // 2
        coluna = idx % 2
        
        x = 0.7 + (coluna * 4.8)
        y = 1.2 + (linha * 1.5)
        
        # Box com métrica
        shape = slide_metricas.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.5), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        shape.line.color.rgb = cor_primaria
        shape.line.width = Pt(2)
        
        # Texto da métrica
        txBox_met = slide_metricas.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(4.1), Inches(0.5))
        tf_met = txBox_met.text_frame
        p_met = tf_met.paragraphs[0]
        p_met.text = nome
        p_met.font.size = Pt(14)
        p_met.font.color.rgb = cor_texto
        
        # Valor
        txBox_val = slide_metricas.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.6), Inches(4.1), Inches(0.5))
        tf_val = txBox_val.text_frame
        p_val = tf_val.paragraphs[0]
        p_val.text = str(valor)
        p_val.font.size = Pt(18)
        p_val.font.bold = True
        p_val.font.color.rgb = cor_primaria
    
    # ===== SLIDES DE GRÁFICOS =====
    if graficos_dict:
        for titulo_grafico, fig in graficos_dict.items():
            slide_grafico = prs.slides.add_slide(prs.slide_layouts[6])
            background_g = slide_grafico.background
            fill_g = background_g.fill
            fill_g.solid()
            fill_g.fore_color.rgb = RGBColor(255, 255, 255)
            
            # Adicionar logo no topo
            if os.path.exists(logo_path):
                try:
                    slide_grafico.shapes.add_picture(logo_path, Inches(0.3), Inches(0.1), height=Inches(0.6))
                except:
                    pass
            
            # Título
            txBox_g_titulo = slide_grafico.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            tf_g_titulo = txBox_g_titulo.text_frame
            p_g_titulo = tf_g_titulo.paragraphs[0]
            p_g_titulo.text = titulo_grafico
            p_g_titulo.font.size = Pt(36)
            p_g_titulo.font.bold = True
            p_g_titulo.font.color.rgb = cor_titulo
            
            # Converter gráfico Plotly para imagem
            try:
                import plotly.io as pio
                
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    tmp_path = tmp.name
                
                # Tentar usar kaleido primeiro
                try:
                    pio.write_image(fig, tmp_path, width=1200, height=700, format='png', engine='kaleido')
                except Exception:
                    # Se kaleido falhar, tentar com outros engines
                    try:
                        pio.write_image(fig, tmp_path, width=1200, height=700, format='png', engine='orca')
                    except Exception:
                        # Se tudo falhar, criar um slide com mensagem de aviso
                        import streamlit as st
                        st.info(f"ℹ️ Gráfico '{titulo_grafico}' não pôde ser convertido - funcionalidade requer Chrome/Chromium instalado")
                        continue
                
                # Verificar se arquivo foi criado com sucesso
                if os.path.exists(tmp_path) and os.path.getsize(tmp_path) > 0:
                    # Inserir imagem no slide
                    slide_grafico.shapes.add_picture(tmp_path, Inches(0.5), Inches(1.1), width=Inches(9))
                    # Limpar arquivo temporário
                    os.unlink(tmp_path)
            except Exception as e:
                import streamlit as st
                st.info(f"ℹ️ Gráfico '{titulo_grafico}' não incluído no relatório - requer Chrome/Chromium instalado no servidor")
    
    # ===== SLIDES DE TOPS COM TABELAS =====
    for titulo_top, df_top in tops_dict.items():
        slide_top = prs.slides.add_slide(prs.slide_layouts[6])
        background_top = slide_top.background
        fill_top = background_top.fill
        fill_top.solid()
        fill_top.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Adicionar logo no topo
        if os.path.exists(logo_path):
            try:
                slide_top.shapes.add_picture(logo_path, Inches(0.3), Inches(0.1), height=Inches(0.6))
            except:
                pass
        
        # Título
        txBox_top_titulo = slide_top.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf_top_titulo = txBox_top_titulo.text_frame
        p_top_titulo = tf_top_titulo.paragraphs[0]
        p_top_titulo.text = titulo_top
        p_top_titulo.font.size = Pt(36)
        p_top_titulo.font.bold = True
        p_top_titulo.font.color.rgb = cor_titulo
        
        # Tabela
        rows = len(df_top) + 1
        cols = len(df_top.columns)
        left = Inches(0.7)
        top = Inches(1.2)
        width = Inches(8.6)
        height = Inches(5.5)
        
        table_shape = slide_top.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Cabeçalho
        for col_idx, col_name in enumerate(df_top.columns):
            cell = table_shape.cell(0, col_idx)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = cor_primaria
            
            # Texto do cabeçalho
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.size = Pt(11)
        
        # Dados
        for row_idx, (_, row) in enumerate(df_top.iterrows(), 1):
            for col_idx, value in enumerate(row):
                cell = table_shape.cell(row_idx, col_idx)
                cell.text = str(value)
                
                # Alternância de cores
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = cor_fundo_claro
                
                # Texto
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = cor_texto
    
    # Converter para bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()

# ==============================
# FILTROS GLOBAIS
# ==============================
def exibir_filtros_globais(df_vendas_original, col_cliente, col_produto, col_vendedor, col_linha, col_data, 
                           col_diretor=None, col_gerente=None, col_gerente_regional=None, 
                           col_supervisor=None, col_coordenador=None, col_consultor=None):
    """
    Exibe painel de filtros globais na sidebar que afetam TODAS as páginas.
    Começa VAZIO - usuário seleciona o que quer
    
    Args:
        df_vendas_original: DataFrame original sem filtros
        col_cliente: Nome da coluna de cliente
        col_produto: Nome da coluna de produto
        col_vendedor: Nome da coluna de vendedor
        col_linha: Nome da coluna de linha
        col_data: Nome da coluna de data
        col_diretor: Nome da coluna de diretor (opcional)
        col_gerente: Nome da coluna de gerente (opcional)
        col_gerente_regional: Nome da coluna de gerente regional (opcional)
        col_supervisor: Nome da coluna de supervisor (opcional)
        col_coordenador: Nome da coluna de coordenador (opcional)
        col_consultor: Nome da coluna de consultor (opcional)
    
    Returns:
        dict: Dicionário com os filtros aplicados
    """
    st.sidebar.markdown("### 🔍 Filtros Globais")
    
    with st.sidebar.expander("📋 Gerenciar Filtros", expanded=True):
        
        # Inicializar filtros em session_state se não existirem
        if 'filtros_globais' not in st.session_state:
            st.session_state.filtros_globais = {
                'produtos': [],
                'vendedores': [],
                'linhas': [],
                'clientes': [],
                'diretores': [],
                'gerentes': [],
                'gerentes_regionais': [],
                'supervisores': [],
                'coordenadores': [],
                'consultores': [],
                'data_inicio': None,
                'data_fim': None
            }
        
        filtros = st.session_state.filtros_globais
        
        # -------- PRODUTO --------
        st.markdown("**📦 Produtos**")
        produtos = sorted(df_vendas_original[col_produto].dropna().unique().tolist())
        
        filtros['produtos'] = st.multiselect(
            "Selecione Produtos:",
            options=produtos,
            default=filtros.get('produtos', []),
            key='filtro_produtos'
        )
        if filtros['produtos']:
            st.caption(f"✅ {len(filtros['produtos'])} produto(s)")
        
        # -------- VENDEDOR --------
        st.markdown("**👤 Vendedores**")
        vendedores = sorted(df_vendas_original[col_vendedor].dropna().unique().tolist())
        
        filtros['vendedores'] = st.multiselect(
            "Selecione Vendedores:",
            options=vendedores,
            default=filtros.get('vendedores', []),
            key='filtro_vendedores'
        )
        if filtros['vendedores']:
            st.caption(f"✅ {len(filtros['vendedores'])} vendedor(es)")
        
        # -------- LINHA --------
        st.markdown("**🏢 Linhas**")
        linhas = sorted(df_vendas_original[col_linha].dropna().unique().tolist())
        
        filtros['linhas'] = st.multiselect(
            "Selecione Linhas:",
            options=linhas,
            default=filtros.get('linhas', []),
            key='filtro_linhas'
        )
        if filtros['linhas']:
            st.caption(f"✅ {len(filtros['linhas'])} linha(s)")
        
        # -------- CLIENTE --------
        st.markdown("**🤝 Clientes**")
        clientes = sorted(df_vendas_original[col_cliente].dropna().unique().tolist())
        
        filtros['clientes'] = st.multiselect(
            "Selecione Clientes:",
            options=clientes,
            default=filtros.get('clientes', []),
            key='filtro_clientes'
        )
        if filtros['clientes']:
            st.caption(f"✅ {len(filtros['clientes'])} cliente(s)")
        
        # -------- HIERARQUIA --------
        st.markdown("---")
        st.markdown("**👔 Hierarquia**")
        
        # Diretor
        if col_diretor and col_diretor != 'Nenhuma':
            diretores = sorted(df_vendas_original[col_diretor].dropna().unique().tolist())
            if diretores:
                filtros['diretores'] = st.multiselect(
                    "Diretores:",
                    options=diretores,
                    default=filtros.get('diretores', []),
                    key='filtro_diretores'
                )
                if filtros['diretores']:
                    st.caption(f"✅ {len(filtros['diretores'])} diretor(es)")
        
        # Gerente Regional
        if col_gerente_regional and col_gerente_regional != 'Nenhuma':
            gerentes_regionais = sorted(df_vendas_original[col_gerente_regional].dropna().unique().tolist())
            if gerentes_regionais:
                filtros['gerentes_regionais'] = st.multiselect(
                    "Gerentes Regionais:",
                    options=gerentes_regionais,
                    default=filtros.get('gerentes_regionais', []),
                    key='filtro_gerentes_regionais'
                )
                if filtros['gerentes_regionais']:
                    st.caption(f"✅ {len(filtros['gerentes_regionais'])} ger. regional(is)")
        
        # Gerente
        if col_gerente and col_gerente != 'Nenhuma':
            gerentes = sorted(df_vendas_original[col_gerente].dropna().unique().tolist())
            if gerentes:
                filtros['gerentes'] = st.multiselect(
                    "Gerentes:",
                    options=gerentes,
                    default=filtros.get('gerentes', []),
                    key='filtro_gerentes'
                )
                if filtros['gerentes']:
                    st.caption(f"✅ {len(filtros['gerentes'])} gerente(s)")
        
        # Supervisor
        if col_supervisor and col_supervisor != 'Nenhuma':
            supervisores = sorted(df_vendas_original[col_supervisor].dropna().unique().tolist())
            if supervisores:
                filtros['supervisores'] = st.multiselect(
                    "Supervisores:",
                    options=supervisores,
                    default=filtros.get('supervisores', []),
                    key='filtro_supervisores'
                )
                if filtros['supervisores']:
                    st.caption(f"✅ {len(filtros['supervisores'])} supervisor(es)")
        
        # Coordenador
        if col_coordenador and col_coordenador != 'Nenhuma':
            coordenadores = sorted(df_vendas_original[col_coordenador].dropna().unique().tolist())
            if coordenadores:
                filtros['coordenadores'] = st.multiselect(
                    "Coordenadores:",
                    options=coordenadores,
                    default=filtros.get('coordenadores', []),
                    key='filtro_coordenadores'
                )
                if filtros['coordenadores']:
                    st.caption(f"✅ {len(filtros['coordenadores'])} coordenador(es)")
        
        # Consultor
        if col_consultor and col_consultor != 'Nenhuma':
            consultores = sorted(df_vendas_original[col_consultor].dropna().unique().tolist())
            if consultores:
                filtros['consultores'] = st.multiselect(
                    "Consultores:",
                    options=consultores,
                    default=filtros.get('consultores', []),
                    key='filtro_consultores'
                )
                if filtros['consultores']:
                    st.caption(f"✅ {len(filtros['consultores'])} consultor(es)")
        
        # -------- DATA RANGE --------
        st.markdown("**📅 Período de Datas**")
        df_temp = df_vendas_original.copy()
        df_temp[col_data] = pd.to_datetime(df_temp[col_data])
        data_min = df_temp[col_data].min()
        data_max = df_temp[col_data].max()
        
        # Verificar se as datas são válidas
        if pd.isna(data_min) or pd.isna(data_max):
            st.warning("⚠️ Datas inválidas ou ausentes nos dados")
            filtros['data_inicio'] = None
            filtros['data_fim'] = None
        else:
            data_min = data_min.date()
            data_max = data_max.date()
            
            # Se não tiver datas selecionadas, começa vazio (None)
            default_data_inicio = filtros.get('data_inicio')
            default_data_fim = filtros.get('data_fim')
            
            if default_data_inicio is None or default_data_fim is None:
                default_range = (data_min, data_max)
            else:
                default_range = (default_data_inicio.date(), default_data_fim.date())
            
            data_inicio, data_fim = st.select_slider(
                "Selecione o período:",
                options=pd.date_range(start=data_min, end=data_max, freq='D').date,
                value=default_range,
                key='filtro_datas'
            )
            
            filtros['data_inicio'] = pd.Timestamp(data_inicio)
            filtros['data_fim'] = pd.Timestamp(data_fim)
            
            st.caption(f"📅 Período: {data_inicio} a {data_fim}")
        
        # -------- BOTÕES --------
        st.markdown("---")
        col_btn1, col_btn2 = st.columns(2)
        
        with col_btn1:
            if st.button("🧹 Limpar Filtros", use_container_width=True):
                st.session_state.filtros_globais = {
                    'produtos': [],
                    'vendedores': [],
                    'linhas': [],
                    'clientes': [],
                    'diretores': [],
                    'gerentes': [],
                    'gerentes_regionais': [],
                    'supervisores': [],
                    'coordenadores': [],
                    'consultores': [],
                    'data_inicio': pd.Timestamp(data_min),
                    'data_fim': pd.Timestamp(data_max)
                }
                st.rerun()
        
        with col_btn2:
            if st.button("✅ Aplicar", use_container_width=True):
                st.session_state.filtros_aplicados = True
                st.rerun()
        
        # -------- INFO --------
        total_registros = len(df_vendas_original)
        filtros_ativos = sum([
            len(filtros.get('produtos', [])),
            len(filtros.get('vendedores', [])),
            len(filtros.get('linhas', [])),
            len(filtros.get('clientes', [])),
            len(filtros.get('diretores', [])),
            len(filtros.get('gerentes', [])),
            len(filtros.get('gerentes_regionais', [])),
            len(filtros.get('supervisores', [])),
            len(filtros.get('coordenadores', [])),
            len(filtros.get('consultores', []))
        ])
        
        st.markdown("---")
        st.info(f"""
        📊 **Status dos Filtros:**
        - Total registros: {total_registros:,}
        - Filtros ativos: {filtros_ativos}
        """)
    
    return filtros

def aplicar_filtros_globais(df_original, filtros, col_cliente, col_produto, col_vendedor, col_linha, col_data,
                           col_diretor=None, col_gerente=None, col_gerente_regional=None, 
                           col_supervisor=None, col_coordenador=None, col_consultor=None):
    """
    Aplica os filtros globais ao DataFrame.
    Se nenhum filtro for selecionado, retorna todos os dados.
    
    Args:
        df_original: DataFrame original
        filtros: Dicionário de filtros retornado por exibir_filtros_globais()
        col_cliente: Nome da coluna de cliente
        col_produto: Nome da coluna de produto
        col_vendedor: Nome da coluna de vendedor
        col_linha: Nome da coluna de linha
        col_data: Nome da coluna de data
        col_diretor: Nome da coluna de diretor (opcional)
        col_gerente: Nome da coluna de gerente (opcional)
        col_gerente_regional: Nome da coluna de gerente regional (opcional)
        col_supervisor: Nome da coluna de supervisor (opcional)
        col_coordenador: Nome da coluna de coordenador (opcional)
        col_consultor: Nome da coluna de consultor (opcional)
    
    Returns:
        DataFrame: DataFrame filtrado (ou original se sem filtros)
    """
    df_filtrado = df_original.copy()
    
    # Se nenhum filtro foi selecionado, retorna todos os dados
    tem_filtro = any([
        filtros.get('produtos'),
        filtros.get('vendedores'),
        filtros.get('linhas'),
        filtros.get('clientes'),
        filtros.get('diretores'),
        filtros.get('gerentes'),
        filtros.get('gerentes_regionais'),
        filtros.get('supervisores'),
        filtros.get('coordenadores'),
        filtros.get('consultores')
    ])
    
    if not tem_filtro:
        # Apenas aplicar date range se selecionado diferente do padrão
        if filtros.get('data_inicio') and filtros.get('data_fim'):
            df_temp = df_filtrado.copy()
            df_temp[col_data] = pd.to_datetime(df_temp[col_data], errors='coerce')
            
            # Remover linhas com datas inválidas
            df_filtrado = df_filtrado[df_temp[col_data].notna()]
            df_temp = df_temp[df_temp[col_data].notna()]
            
            if not df_temp.empty:
                df_filtrado = df_filtrado[
                    (df_temp[col_data] >= filtros['data_inicio']) & 
                    (df_temp[col_data] <= filtros['data_fim'])
                ]
        return df_filtrado
    
    # Aplicar filtros básicos
    if filtros.get('produtos'):
        df_filtrado = df_filtrado[df_filtrado[col_produto].isin(filtros['produtos'])]
    
    if filtros.get('vendedores'):
        df_filtrado = df_filtrado[df_filtrado[col_vendedor].isin(filtros['vendedores'])]
    
    if filtros.get('linhas'):
        df_filtrado = df_filtrado[df_filtrado[col_linha].isin(filtros['linhas'])]
    
    if filtros.get('clientes'):
        df_filtrado = df_filtrado[df_filtrado[col_cliente].isin(filtros['clientes'])]
    
    # Aplicar filtros de hierarquia
    if filtros.get('diretores') and col_diretor and col_diretor != 'Nenhuma':
        df_filtrado = df_filtrado[df_filtrado[col_diretor].isin(filtros['diretores'])]
    
    if filtros.get('gerentes') and col_gerente and col_gerente != 'Nenhuma':
        df_filtrado = df_filtrado[df_filtrado[col_gerente].isin(filtros['gerentes'])]
    
    if filtros.get('gerentes_regionais') and col_gerente_regional and col_gerente_regional != 'Nenhuma':
        df_filtrado = df_filtrado[df_filtrado[col_gerente_regional].isin(filtros['gerentes_regionais'])]
    
    if filtros.get('supervisores') and col_supervisor and col_supervisor != 'Nenhuma':
        df_filtrado = df_filtrado[df_filtrado[col_supervisor].isin(filtros['supervisores'])]
    
    if filtros.get('coordenadores') and col_coordenador and col_coordenador != 'Nenhuma':
        df_filtrado = df_filtrado[df_filtrado[col_coordenador].isin(filtros['coordenadores'])]
    
    if filtros.get('consultores') and col_consultor and col_consultor != 'Nenhuma':
        df_filtrado = df_filtrado[df_filtrado[col_consultor].isin(filtros['consultores'])]
    
    # Aplicar filtro de data
    if filtros.get('data_inicio') and filtros.get('data_fim'):
        df_temp = df_filtrado.copy()
        df_temp[col_data] = pd.to_datetime(df_temp[col_data], errors='coerce')
        
        # Remover linhas com datas inválidas antes de filtrar
        df_filtrado = df_filtrado[df_temp[col_data].notna()]
        df_temp = df_temp[df_temp[col_data].notna()]
        
        if not df_temp.empty:
            df_filtrado = df_filtrado[
                (df_temp[col_data] >= filtros['data_inicio']) & 
                (df_temp[col_data] <= filtros['data_fim'])
            ]
    
    return df_filtrado

def obter_dados_filtrados():
    """
    Função auxiliar para páginas internas obterem dados filtrados dos filtros globais.
    Retorna df_vendas e df_devolucoes já filtrados e prontos para usar.
    
    Returns:
        tuple: (df_vendas_filtrado, df_devolucoes_filtrado)
    """
    # Os dados já vêm filtrados do session_state (aplicados em app.py)
    df_vendas = st.session_state.get('df_vendas', pd.DataFrame())
    df_devolucoes = st.session_state.get('df_devolucoes', pd.DataFrame())
    
    return df_vendas, df_devolucoes